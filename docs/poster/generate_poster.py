"""
Generate UniSchedule AI A1 landscape poster — visual design, ~310 body words.
Run: python docs/poster/generate_poster.py
"""

import re
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Cm, Pt

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY = RGBColor(0x14, 0x1F, 0x4A)
NAVY_DEEP = RGBColor(0x0B, 0x14, 0x36)
# Supervisor change: green accent replaced with light blue across the poster.
# Names kept as GREEN/GREEN_DARK to avoid touching every call site; values are now blue.
GREEN = RGBColor(0x3D, 0x9B, 0xE0)
GREEN_DARK = RGBColor(0x1F, 0x6F, 0xC0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF6, 0xF7, 0xFB)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE2, 0xE5, 0xED)
SHADOW = RGBColor(0xD8, 0xDB, 0xE3)
TEXT = RGBColor(0x1F, 0x23, 0x2E)
MUTED = RGBColor(0x55, 0x5B, 0x6B)
LIGHT_BLUE = RGBColor(0xE8, 0xED, 0xF8)
SOFT_GREEN = RGBColor(0xE6, 0xF0, 0xFB)
AMBER = RGBColor(0xF5, 0xA6, 0x23)
LIGHT_AMBER = RGBColor(0xFD, 0xF1, 0xDE)
GOLD = RGBColor(0xC8, 0x9B, 0x00)
LIGHT_GOLD = RGBColor(0xFC, 0xF3, 0xCF)
TITLE_SUB = RGBColor(0xCC, 0xDD, 0xFF)
TITLE_META = RGBColor(0xAA, 0xBB, 0xDD)
TITLE_MOTIF = RGBColor(0x1C, 0x2A, 0x5E)
HERO_GREEN = RGBColor(0xE9, 0xF3, 0xFD)

YOUR_NAME = "[Your Full Name]"
STUDENT_ID = "[Student ID]"
SUPERVISOR = "[Supervisor Name]"
OUTPUT = Path(__file__).parent / "UniSchedule-AI-Poster-v10.pptx"

# ── Geometry (cm) ─────────────────────────────────────────────────────────────
SLIDE_W, SLIDE_H = 84.1, 59.4
MARGIN, GUTTER = 2.0, 1.0
COL_W = (SLIDE_W - 2 * MARGIN - 2 * GUTTER) / 3
TITLE_H = 9.0
FOOTER_H = 2.0
SECTION_GAP = 0.6
CARD_HEADER = 1.85
CONTENT_TOP = TITLE_H + 0.6
CONTENT_BOTTOM = SLIDE_H - FOOTER_H - 0.4
CONTENT_H = CONTENT_BOTTOM - CONTENT_TOP
COL_X = [MARGIN, MARGIN + COL_W + GUTTER, MARGIN + 2 * (COL_W + GUTTER)]

# ── Body text (target 300–500 words) ──────────────────────────────────────────
BODY = {
    "contribution": (
        "A hybrid rules-and-NLP scheduler with deterministic academic-constraint "
        "validation, evaluated against manual scheduling."
    ),
    "intro": (
        "Academic meetings — consultations, supervisions, project work — are booked "
        "manually by email, causing clashes, delays, and poor room use."
    ),
}

GAP_BULLETS = [
    "Corporate tools ignore academic roles and exam periods",
    "Timetable systems: fixed semesters, no ad-hoc booking",
    "Literature: strong theory, few deployable systems",
]

NOVEL_BULLETS = [
    "NLP interprets requests; rule engine validates every booking",
    "Enforces exams, capacity, availability, role-based access",
    "Deployable Next.js + Laravel + MySQL web platform",
]

RESULTS_BULLETS = [
    "Zero clash false negatives · API < 3 seconds",
    "SUS ≥ 68 · ≥ 25% faster than manual booking",
    "Evaluated Weeks 7–11 with 5–8 participants",
]

FEATURES = [
    ("01", "Sanctum Auth", "Role-based API access: student, lecturer, admin.", NAVY),
    ("02", "Constraint Engine", "Exams, capacity, availability validated server-side.",
     GREEN_DARK),
    ("03", "NLP Requests", "Natural language parsed, then rule-checked.", NAVY),
    ("04", "Clash Detection", "Time, room, person conflicts flagged instantly.",
     GREEN_DARK),
    ("05", "Room Allocation", "Capacity-aware matching with fallback rooms.", NAVY),
    ("06", "Analytics", "Room usage, peak demand, clash patterns.", GREEN_DARK),
]

OBJECTIVES = [
    ("O1", "Hybrid AI Scheduler", "NLP requests + deterministic rule engine.", GREEN_DARK),
    ("O2", "Academic Constraints", "Exams, capacity, availability enforced.", NAVY),
    ("O3", "Clash Detection", "Time, room, participant conflicts resolved.", GREEN_DARK),
    ("O4", "Empirical Evaluation", "AI-assisted vs manual, measured.", NAVY),
]


def word_count(*texts):
    return len(re.findall(r"\b\w+\b", " ".join(texts)))


# ── Primitives ───────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill, line=None, lw=Pt(1), rounded=False):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    s = slide.shapes.add_shape(shape_type, Cm(l), Cm(t), Cm(w), Cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s


def tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))


def write(shape, lines, anchor=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Cm(0.35)
    tf.margin_right = Cm(0.35)
    tf.margin_top = Cm(0.15)
    tf.margin_bottom = Cm(0.12)
    for i, item in enumerate(lines):
        text, size, bold, color = item[:4]
        align = item[4] if len(item) > 4 else PP_ALIGN.LEFT
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.alignment = align
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(4)
        p.line_spacing = 1.18


# ── Card with shadow, navy header, green accent line, numbered badge ─────────
def card(slide, col, top, height, title, badge=None):
    x = COL_X[col]
    add_rect(slide, x + 0.12, top + 0.14, COL_W, height, SHADOW)
    add_rect(slide, x, top, COL_W, height, CARD_BG, CARD_BORDER, Pt(0.75))
    add_rect(slide, x, top, COL_W, CARD_HEADER, NAVY_DEEP)
    add_rect(slide, x, top + CARD_HEADER, COL_W, 0.14, GREEN)
    if badge:
        write(
            tb(slide, x + 0.45, top + 0.28, 2.4, CARD_HEADER - 0.3),
            [(badge, 22, True, GREEN)],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        write(
            tb(slide, x + 2.9, top + 0.28, COL_W - 3.3, CARD_HEADER - 0.3),
            [(title, 30, True, WHITE)],
            anchor=MSO_ANCHOR.MIDDLE,
        )
    else:
        write(
            tb(slide, x + 0.45, top + 0.28, COL_W - 0.9, CARD_HEADER - 0.3),
            [(title, 30, True, WHITE)],
            anchor=MSO_ANCHOR.MIDDLE,
        )
    body_top = top + CARD_HEADER + 0.5
    body_h = height - CARD_HEADER - 0.65
    return x + 0.45, body_top, COL_W - 0.9, body_h


# ── Section renderers ────────────────────────────────────────────────────────
def draw_rq_badges(slide, x, y, w):
    rqs = [
        ("RQ1", "Scheduling difficulty"),
        ("RQ2", "Clash frequency"),
        ("RQ3", "Rule violations"),
        ("RQ4", "Coordination time"),
        ("RQ5", "User satisfaction"),
    ]
    bw = (w - 0.8) / 5
    for i, (label, sub) in enumerate(rqs):
        bx = x + i * (bw + 0.2)
        cx = bx + bw / 2 - 1.1
        circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Cm(cx), Cm(y), Cm(2.2), Cm(2.2)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = GREEN_DARK if i % 2 == 0 else NAVY
        circle.line.color.rgb = WHITE
        circle.line.width = Pt(2)
        write(
            tb(slide, cx, y, 2.2, 2.2),
            [(label, 20, True, WHITE, PP_ALIGN.CENTER)],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        write(
            tb(slide, bx - 0.1, y + 2.35, bw + 0.2, 1.8),
            [(sub, 14, False, MUTED, PP_ALIGN.CENTER)],
        )


def draw_gap_triad(slide, x, y, w, h):
    items = [
        ("Corporate Tools", "Business-focused;\nno academic rules", LIGHT_BLUE, NAVY),
        ("Timetable Systems", "Fixed semester only;\nno ad-hoc support", SOFT_GREEN, GREEN_DARK),
        ("Academic Literature", "Strong theory;\nfew deployable systems", LIGHT_AMBER, AMBER),
    ]
    bw = (w - 0.6) / 3
    for i, (title, sub, bg, accent) in enumerate(items):
        bx = x + i * (bw + 0.3)
        add_rect(slide, bx, y, bw, h, bg)
        add_rect(slide, bx, y, bw, 0.25, accent)
        write(
            tb(slide, bx + 0.2, y + 0.45, bw - 0.4, 1.0),
            [(title, 13, True, NAVY, PP_ALIGN.CENTER)],
        )
        write(
            tb(slide, bx + 0.2, y + 1.55, bw - 0.4, h - 1.7),
            [(sub, 11, False, MUTED, PP_ALIGN.CENTER)],
        )


def draw_architecture(slide, x, y, w, h):
    cx = x + w / 2
    layers = [
        ("Next.js Frontend", "Student & lecturer UI · TypeScript", NAVY, WHITE, "HTTPS"),
        ("Laravel REST API", "Sanctum · Scheduler · Clash · NLP · Rooms",
         GREEN_DARK, WHITE, "Eloquent / SQL"),
        ("MySQL Database", "Users · Meetings · Rooms · Constraints",
         RGBColor(0xD2, 0xD7, 0xE3), TEXT, None),
    ]
    n = len(layers)
    gap = 1.3
    bh = (h - gap * (n - 1) - 0.3) / n
    bx = x + w * 0.05
    bw = w * 0.90
    cy = y + 0.15
    for i, (title, sub, bg, fg, arrow_label) in enumerate(layers):
        add_rect(slide, bx + 0.1, cy + 0.12, bw, bh, SHADOW, rounded=True)
        add_rect(slide, bx, cy, bw, bh, bg, rounded=True)
        write(
            tb(slide, bx, cy + bh * 0.14, bw, bh * 0.46),
            [(title, 23, True, fg, PP_ALIGN.CENTER)],
        )
        sub_color = WHITE if bg in (NAVY, GREEN_DARK) else MUTED
        write(
            tb(slide, bx, cy + bh * 0.60, bw, bh * 0.38),
            [(sub, 15, False, sub_color, PP_ALIGN.CENTER)],
        )
        if i < n - 1:
            ay = cy + bh + 0.12
            arr = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
                Cm(cx - 0.55), Cm(ay), Cm(1.1), Cm(gap - 0.28),
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = GREEN
            arr.line.fill.background()
            if arrow_label:
                write(
                    tb(slide, cx + 0.85, ay - 0.05, w * 0.42, gap),
                    [(arrow_label, 14, True, MUTED)],
                    anchor=MSO_ANCHOR.MIDDLE,
                )
        cy += bh + gap


def draw_timeline(slide, x, y, w, h):
    steps = [
        ("1", "Literature", "Review"),
        ("2", "Requirements", "Survey + interviews"),
        ("3", "Design", "Architecture & ER"),
        ("4", "Build", "Next.js · Laravel"),
        ("5", "Evaluate", "Manual vs AI"),
    ]
    n = len(steps)
    sw = w / n
    line_y = y + 0.7
    add_rect(slide, x + sw * 0.3, line_y + 0.65, w - sw * 0.6, 0.16, GREEN)
    for i, (num, t1, t2) in enumerate(steps):
        sx = x + i * sw + sw / 2 - 0.7
        circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Cm(sx), Cm(line_y), Cm(1.4), Cm(1.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = NAVY if i % 2 == 0 else GREEN_DARK
        circle.line.color.rgb = WHITE
        circle.line.width = Pt(2)
        write(
            tb(slide, sx, line_y, 1.4, 1.4),
            [(num, 20, True, WHITE, PP_ALIGN.CENTER)],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        write(
            tb(slide, x + i * sw, line_y + 1.7, sw, 2.2),
            [(t1, 17, True, NAVY, PP_ALIGN.CENTER),
             (t2, 13, False, MUTED, PP_ALIGN.CENTER)],
        )


def _style_data_labels(plot, size=12):
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(size)
    dl.font.bold = True
    dl.font.name = "Arial"
    dl.number_format = "0"
    dl.number_format_is_linked = False


def add_moscow_chart(slide, x, y, w, h):
    data = CategoryChartData()
    data.categories = ["Must", "Should", "Could", "Won't"]
    data.add_series("FR count", (9, 3, 2, 3))
    cf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Cm(x), Cm(y), Cm(w), Cm(h), data
    )
    chart = cf.chart
    chart.has_legend = False
    chart.chart_title.has_text_frame = False
    colors = [NAVY, GREEN_DARK, AMBER, RGBColor(0xBB, 0xBB, 0xBB)]
    for i, c in enumerate(colors):
        pt = chart.plots[0].series[0].points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = c
    plot = chart.plots[0]
    _style_data_labels(plot, 18)
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.gap_width = 80
    chart.category_axis.tick_labels.font.size = Pt(17)
    chart.category_axis.tick_labels.font.bold = True
    chart.category_axis.format.line.color.rgb = CARD_BORDER
    chart.value_axis.visible = False
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.maximum_scale = 11


def add_targets_chart(slide, x, y, w, h):
    data = CategoryChartData()
    data.categories = ["SUS (\u226568)", "Time saved %", "Scenarios", "API (sec)"]
    data.add_series("Target", (68, 25, 20, 3))
    cf = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Cm(x), Cm(y), Cm(w), Cm(h), data
    )
    chart = cf.chart
    chart.has_legend = False
    chart.chart_title.has_text_frame = False
    bar_colors = [NAVY, GREEN_DARK, AMBER, GREEN_DARK]
    for i, c in enumerate(bar_colors):
        pt = chart.plots[0].series[0].points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = c
    plot = chart.plots[0]
    _style_data_labels(plot, 16)
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.gap_width = 60
    chart.category_axis.tick_labels.font.size = Pt(15)
    chart.category_axis.tick_labels.font.bold = True
    chart.category_axis.format.line.color.rgb = CARD_BORDER
    chart.value_axis.visible = False
    chart.value_axis.has_major_gridlines = False


def draw_feature_cards(slide, x, y, w, h):
    cols, rows = 2, 3
    cw = (w - 0.4) / cols
    ch = (h - 0.6) / rows
    for i, (num, title, sub, accent) in enumerate(FEATURES):
        col_i, row_i = i % cols, i // cols
        fx = x + col_i * (cw + 0.4)
        fy = y + row_i * (ch + 0.3)
        add_rect(slide, fx + 0.08, fy + 0.10, cw, ch, SHADOW)
        add_rect(slide, fx, fy, cw, ch, WHITE, CARD_BORDER, Pt(0.75))
        add_rect(slide, fx, fy, 0.45, ch, accent)
        write(
            tb(slide, fx + 0.7, fy + 0.2, 2.0, 0.7),
            [(num, 16, True, accent)],
        )
        write(
            tb(slide, fx + 0.7, fy + 0.85, cw - 0.9, 0.9),
            [(title, 19, True, NAVY)],
        )
        write(
            tb(slide, fx + 0.7, fy + 1.85, cw - 0.9, ch - 1.95),
            [(sub, 14, False, MUTED)],
        )


def draw_title_motif(slide):
    for i in range(1, 14):
        lx = i * (SLIDE_W / 14)
        add_rect(slide, lx, 0, 0.045, TITLE_H, TITLE_MOTIF)
    for r in (5.0, 8.5, 12.0):
        c = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Cm(SLIDE_W - 12 - r / 2), Cm(TITLE_H / 2 - r / 2), Cm(r), Cm(r),
        )
        c.fill.background()
        c.line.color.rgb = TITLE_MOTIF
        c.line.width = Pt(2)


def draw_references(slide, x, y, w, h):
    refs = [
        ("Burke, E.K. & Petrovic, S. (2002) 'Recent research directions in automated "
         "timetabling', Eur. J. Operational Research, 140(2)."),
        ("Müller, T., Rudová, H. & Müllerová, Z. (2025) 'Real-world university course "
         "timetabling', Journal of Scheduling, 28(2)."),
        ("Schaerf, A. (1999) 'A survey of automated timetabling', Artificial "
         "Intelligence Review, 13(2)."),
        ("Wijerathne, O. et al. (2025) 'ScheduleMe', arXiv:2509.25693."),
    ]
    lines = [(f"▸  {r}", 14, False, TEXT) for r in refs]
    write(tb(slide, x, y, w, h), lines)


def draw_workplan(slide, x, y, w, h):
    phases = [
        ("Literature review", 1, 3, NAVY),
        ("Requirements + survey", 3, 5, GREEN_DARK),
        ("System design", 4, 6, NAVY),
        ("Development", 5, 9, GREEN_DARK),
        ("Evaluation / UAT", 9, 11, NAVY),
        ("Dissertation write-up", 2, 12, AMBER),
    ]
    total_weeks = 12
    label_w = w * 0.38
    track_x = x + label_w
    track_w = w - label_w
    top = y + 1.1
    avail = h - 1.5
    n = len(phases)
    row_h = avail / n
    bar_h = row_h * 0.62

    for wk in range(2, total_weeks + 1, 2):
        lx = track_x + track_w * (wk / total_weeks)
        write(
            tb(slide, lx - 0.7, y - 0.05, 1.4, 0.8),
            [(f"W{wk}", 12, True, MUTED, PP_ALIGN.CENTER)],
        )

    for i, (label, s, e, color) in enumerate(phases):
        ry = top + i * row_h
        write(
            tb(slide, x, ry, label_w - 0.2, row_h),
            [(label, 14, True, TEXT)],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        bx = track_x + track_w * (s / total_weeks)
        bw = track_w * ((e - s) / total_weeks)
        add_rect(slide, bx, ry + (row_h - bar_h) / 2, bw, bar_h, color, rounded=True)

    cur = 5
    mx = track_x + track_w * (cur / total_weeks)
    marker = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(mx - 0.04), Cm(top), Cm(0.08), Cm(avail)
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = AMBER
    marker.line.fill.background()
    write(
        tb(slide, mx - 1.0, top + avail + 0.02, 2.0, 0.7),
        [("now", 12, True, AMBER, PP_ALIGN.CENTER)],
    )


def draw_risks(slide, x, y, w, h):
    items = [
        ("Low survey participation",
         "University channels, supervisor support, peer outreach"),
        ("LLM misinterpretation",
         "Deterministic constraint validation before any booking"),
        ("Small sample (5–8)",
         "Within-subject A/B design improves statistical power"),
        ("Scope creep",
         "MoSCoW priorities; voice & ML deferred to future work"),
    ]
    lines = []
    for risk, mit in items:
        lines.append((f"⚠  {risk}", 11, True, NAVY))
        lines.append((f"     → {mit}", 9, False, MUTED))
    lines.append(("Limitation: single-institution scope; no live LMS integration.",
                  9, False, TEXT))
    write(tb(slide, x, y, w, h), lines)


def draw_tiles(slide, x, y, w, h, items, cols, rows):
    cw = (w - 0.4 * (cols - 1)) / cols
    ch = (h - 0.3 * (rows - 1)) / rows
    for i, (num, title, sub, accent) in enumerate(items):
        col_i, row_i = i % cols, i // cols
        fx = x + col_i * (cw + 0.4)
        fy = y + row_i * (ch + 0.3)
        add_rect(slide, fx + 0.08, fy + 0.10, cw, ch, SHADOW)
        add_rect(slide, fx, fy, cw, ch, WHITE, CARD_BORDER, Pt(0.75))
        add_rect(slide, fx, fy, 0.45, ch, accent)
        write(tb(slide, fx + 0.7, fy + 0.25, 2.0, 0.8), [(num, 18, True, accent)])
        write(tb(slide, fx + 0.7, fy + 1.0, cw - 0.9, 0.9), [(title, 20, True, NAVY)])
        write(tb(slide, fx + 0.7, fy + 2.0, cw - 0.9, ch - 2.1),
              [(sub, 15, False, MUTED)])


def _er_entity(slide, x, y, w, title, fields):
    header_h = 1.15
    row_h = 0.82
    h = header_h + row_h * len(fields)
    add_rect(slide, x + 0.08, y + 0.1, w, h, SHADOW, rounded=True)
    add_rect(slide, x, y, w, h, WHITE, CARD_BORDER, Pt(1.0), rounded=True)
    add_rect(slide, x, y, w, header_h, NAVY)
    write(
        tb(slide, x, y + 0.05, w, header_h - 0.1),
        [(title, 16, True, WHITE, PP_ALIGN.CENTER)],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    for i, f in enumerate(fields):
        fy = y + header_h + i * row_h
        write(
            tb(slide, x + 0.3, fy, w - 0.45, row_h),
            [(f, 13, False, TEXT)],
            anchor=MSO_ANCHOR.MIDDLE,
        )
    return (x, y, w, h)


def _anchor(p, side):
    x, y, w, h = p
    return {
        "l": (x, y + h / 2), "r": (x + w, y + h / 2),
        "t": (x + w / 2, y), "b": (x + w / 2, y + h),
    }[side]


def _er_connect(slide, p_from, p_to, label, frm="r", to="l"):
    x1, y1 = _anchor(p_from, frm)
    x2, y2 = _anchor(p_to, to)
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2)
    )
    conn.line.color.rgb = GREEN_DARK
    conn.line.width = Pt(2.0)
    mx, my = (x1 + x2) / 2, (y1 + y2) / 2
    add_rect(slide, mx - 1.6, my - 0.45, 3.2, 0.9, WHITE)
    write(
        tb(slide, mx - 1.6, my - 0.45, 3.2, 0.9),
        [(label, 12, True, GREEN_DARK, PP_ALIGN.CENTER)],
        anchor=MSO_ANCHOR.MIDDLE,
    )


def draw_er_diagram(slide, x, y, w, h):
    ebw = 7.4
    colA = x
    colB = x + (w - ebw) / 2
    colC = x + w - ebw

    users = _er_entity(slide, colA, y + 0.8, ebw, "users",
                       ["PK  id", "name", "email", "role"])
    avail = _er_entity(slide, colA, y + 6.4, ebw, "availabilities",
                       ["PK  id", "FK  user_id", "day, start–end"])
    meetings = _er_entity(slide, colB, y + 3.2, ebw, "meetings",
                          ["PK  id", "FK  organizer_id", "FK  room_id",
                           "start–end, status"])
    parts = _er_entity(slide, colB, y + 10.0, ebw, "meeting_participants",
                       ["PK  id", "FK  meeting_id", "FK  user_id", "response"])
    rooms = _er_entity(slide, colC, y + 0.8, ebw, "rooms",
                       ["PK  id", "name", "capacity"])
    clash = _er_entity(slide, colC, y + 5.6, ebw, "clash_records",
                       ["PK  id", "FK  meeting_id", "clash_type"])
    rules = _er_entity(slide, colC, y + 10.4, ebw, "constraint_rules",
                       ["PK  id", "rule_type", "valid_from–to"])

    _er_connect(slide, users, meetings, "organises 1..*", "r", "l")
    _er_connect(slide, users, avail, "has 1..*", "b", "t")
    _er_connect(slide, meetings, parts, "has 1..*", "b", "t")
    _er_connect(slide, parts, users, "joins *..1", "l", "b")
    _er_connect(slide, meetings, rooms, "in *..1", "r", "l")
    _er_connect(slide, meetings, clash, "raises 1..*", "r", "l")
    _er_connect(slide, rules, meetings, "validates", "l", "r")


# ── Main ─────────────────────────────────────────────────────────────────────
def build_poster():
    extra = GAP_BULLETS + NOVEL_BULLETS + RESULTS_BULLETS
    extra += [t + " " + s for _, t, s, _ in FEATURES + OBJECTIVES]
    wc = word_count(*BODY.values(), *extra)
    print(f"Body word count: {wc} (concise/bullet layout for A1 readability)")

    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)

    # ── Title bar ────────────────────────────────────────────────────────────
    add_rect(slide, 0, 0, SLIDE_W, TITLE_H, NAVY_DEEP)
    draw_title_motif(slide)
    add_rect(slide, 0, TITLE_H - 0.2, SLIDE_W, 0.2, GREEN)
    add_rect(slide, MARGIN, 1.35, 7.4, 1.2, GREEN, rounded=True)
    write(
        tb(slide, MARGIN, 1.4, 7.4, 1.1),
        [("FINAL YEAR PROJECT  ·  2026", 15, True, WHITE, PP_ALIGN.CENTER)],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    write(
        tb(slide, MARGIN, 2.55, SLIDE_W - 2 * MARGIN, 3.4),
        [("UniSchedule AI", 78, True, WHITE)],
    )
    write(
        tb(slide, MARGIN, 6.2, SLIDE_W - 2 * MARGIN, 1.4),
        [("Intelligent University Meeting Scheduler — academic-aware, NLP-driven, "
          "real-time clash detection",
          27, False, TITLE_SUB)],
    )
    write(
        tb(slide, MARGIN, 7.65, SLIDE_W - 2 * MARGIN, 1.2),
        [(f"{YOUR_NAME}   ·   {STUDENT_ID}   ·   BEng Computer Science   ·   "
          f"University of Canterbury",
          18, False, TITLE_META)],
    )

    # ── Column 1 — PROBLEM ───────────────────────────────────────────────────
    c1 = [20.5, 12.0, 13.6]
    y = CONTENT_TOP

    x, cy, cw, ch = card(slide, 0, y, c1[0], "Introduction & Objectives", "01")
    write(tb(slide, x, cy, cw, ch * 0.13), [(BODY["intro"], 18, False, TEXT)])
    write(tb(slide, x, cy + ch * 0.15, cw, 0.9), [("Objectives", 22, True, GREEN_DARK)])
    draw_tiles(slide, x, cy + ch * 0.22, cw, ch * 0.42, OBJECTIVES, 2, 2)
    write(tb(slide, x, cy + ch * 0.67, cw, 0.9),
          [("Research Questions", 22, True, GREEN_DARK)])
    draw_rq_badges(slide, x, cy + ch * 0.77, cw)
    y += c1[0] + SECTION_GAP

    x, cy, cw, ch = card(slide, 0, y, c1[1], "Research Gap", "02")
    write(tb(slide, x, cy, cw, ch * 0.48),
          [(f"•  {b}", 18, False, TEXT) for b in GAP_BULLETS])
    gold_y = cy + ch * 0.52
    gold_h = ch * 0.46
    add_rect(slide, x - 0.15, gold_y, cw + 0.3, gold_h, LIGHT_GOLD, rounded=True)
    add_rect(slide, x - 0.15, gold_y, 0.22, gold_h, GOLD)
    write(
        tb(slide, x + 0.15, gold_y + 0.1, cw - 0.2, gold_h - 0.2),
        [
            ("WHY UNISCHEDULE AI IS UNIQUE", 15, True, GOLD),
            ("Academic-constraint validation + NLP + real-time clash detection — "
             "empirically evaluated.", 18, True, NAVY),
        ],
    )
    y += c1[1] + SECTION_GAP

    x, cy, cw, ch = card(slide, 0, y, c1[2], "Novel Approach & Contribution", "03")
    hero_h = ch * 0.34
    add_rect(slide, x - 0.15, cy - 0.1, cw + 0.3, hero_h, HERO_GREEN, rounded=True)
    add_rect(slide, x - 0.15, cy - 0.1, 0.22, hero_h, GREEN_DARK)
    write(
        tb(slide, x + 0.15, cy, cw - 0.2, hero_h),
        [
            ("KEY CONTRIBUTION", 15, True, GREEN_DARK),
            (BODY["contribution"], 18, True, NAVY),
        ],
    )
    write(
        tb(slide, x, cy + hero_h + 0.25, cw, ch * 0.32),
        [(f"•  {b}", 16, False, TEXT) for b in NOVEL_BULLETS],
    )
    pillars = [
        ("Hybrid AI", GREEN_DARK),
        ("Constraint-aware", NAVY),
        ("Academic rules", AMBER),
        ("Deployable web", GREEN_DARK),
    ]
    pw = (cw - 0.45) / 4
    py = cy + ch * 0.79
    for i, (label, c) in enumerate(pillars):
        px = x + i * (pw + 0.15)
        add_rect(slide, px, py, pw, 2.0, c, rounded=True)
        write(
            tb(slide, px, py, pw, 2.0),
            [(label, 14, True, WHITE, PP_ALIGN.CENTER)],
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # ── Column 2 — APPROACH ──────────────────────────────────────────────────
    c2 = [11.5, 16.5, 8.0, 8.5]
    y = CONTENT_TOP

    x, cy, cw, ch = card(slide, 1, y, c2[0], "System Architecture", "04")
    draw_architecture(slide, x, cy, cw, ch)
    y += c2[0] + SECTION_GAP

    x, cy, cw, ch = card(slide, 1, y, c2[1], "ER Diagram — Database Schema", "05")
    draw_er_diagram(slide, x, cy, cw, ch)
    y += c2[1] + SECTION_GAP

    x, cy, cw, ch = card(slide, 1, y, c2[2], "Research Methodology", "06")
    write(
        tb(slide, x, cy, cw, 0.8),
        [("Applied research with experimental evaluation",
          16, True, NAVY, PP_ALIGN.CENTER)],
    )
    draw_timeline(slide, x, cy + 0.9, cw, ch - 0.9)
    y += c2[2] + SECTION_GAP

    x, cy, cw, ch = card(slide, 1, y, c2[3], "Project Scope — MoSCoW", "07")
    write(
        tb(slide, x, cy, cw, 0.8),
        [("Functional requirements by priority",
          15, False, MUTED, PP_ALIGN.CENTER)],
    )
    add_moscow_chart(slide, x, cy + 0.8, cw, ch - 0.8)

    # ── Column 3 — FEATURES, RESULTS, PLAN, REFERENCES ───────────────────────
    c3 = [15.0, 9.0, 10.5, 11.0]
    y = CONTENT_TOP

    x, cy, cw, ch = card(slide, 2, y, c3[0], "Key Features", "08")
    draw_feature_cards(slide, x, cy, cw, ch)
    y += c3[0] + SECTION_GAP

    x, cy, cw, ch = card(slide, 2, y, c3[1], "Expected Results", "09")
    write(tb(slide, x, cy, cw, ch * 0.42),
          [(f"•  {b}", 17, False, TEXT) for b in RESULTS_BULLETS])
    add_targets_chart(slide, x, cy + ch * 0.44, cw, ch * 0.54)
    y += c3[1] + SECTION_GAP

    x, cy, cw, ch = card(slide, 2, y, c3[2], "Plan of Work", "10")
    draw_workplan(slide, x, cy, cw, ch)
    y += c3[2] + SECTION_GAP

    x, cy, cw, ch = card(slide, 2, y, c3[3], "References (Harvard)", "11")
    draw_references(slide, x, cy, cw, ch)

    # ── Footer ───────────────────────────────────────────────────────────────
    fy = SLIDE_H - FOOTER_H
    add_rect(slide, 0, fy, SLIDE_W, FOOTER_H, NAVY_DEEP)
    add_rect(slide, 0, fy, SLIDE_W, 0.14, GREEN)
    write(
        tb(slide, MARGIN, fy + 0.5, (SLIDE_W - 2 * MARGIN) * 0.5, FOOTER_H - 0.6),
        [("University of Canterbury  ·  School of Engineering",
          18, True, WHITE)],
    )
    write(
        tb(slide, SLIDE_W / 2, fy + 0.5, SLIDE_W / 2 - MARGIN, FOOTER_H - 0.6),
        [(f"github.com/bhagya-tharindu/UniSchedule-AI.   ·   Supervisor: {SUPERVISOR}",
          18, False, TITLE_SUB, PP_ALIGN.RIGHT)],
    )

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT.resolve()}")


if __name__ == "__main__":
    build_poster()
